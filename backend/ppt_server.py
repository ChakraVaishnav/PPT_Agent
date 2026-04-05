import os
import shutil
from typing import List

import requests
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt, Emu

from config import Config
from fastmcp import FastMCP

# ─── CONSTANTS ─────────────────────────────────────────────────────────────────
BASE_DIR = r"C:\Users\gunta\Downloads\Calibo"
TEMP_FILE = os.path.join(BASE_DIR, "temp_presentation.pptx")
TEMPLATES_DIR = os.path.join(BASE_DIR, "templates")  # folder with your .pptx templates

mcp = FastMCP("ppt-server")

# ─── THEME REGISTRY ────────────────────────────────────────────────────────────
# Each theme: template file + style overrides applied on top
THEMES = {
    "dark_tech": {
        "template": os.path.join(TEMPLATES_DIR, "dark_tech.pptx"),
        "bg_color": (15, 15, 25),          # near-black
        "title_color": (0, 220, 255),       # cyan
        "body_color": (220, 220, 220),      # light gray
        "title_font": "Arial Black",
        "body_font": "Calibri",
        "title_size": 40,
        "body_size": 18,
    },
    "minimal_light": {
        "template": os.path.join(TEMPLATES_DIR, "minimal_light.pptx"),
        "bg_color": (255, 255, 255),
        "title_color": (30, 30, 30),
        "body_color": (60, 60, 60),
        "title_font": "Georgia",
        "body_font": "Calibri Light",
        "title_size": 38,
        "body_size": 16,
    },
    "corporate_blue": {
        "template": os.path.join(TEMPLATES_DIR, "corporate_blue.pptx"),
        "bg_color": (30, 39, 97),           # deep navy
        "title_color": (255, 255, 255),
        "body_color": (202, 220, 252),      # ice blue
        "title_font": "Calibri",
        "body_font": "Calibri",
        "title_size": 38,
        "body_size": 17,
    },
    "vibrant_creative": {
        "template": os.path.join(TEMPLATES_DIR, "vibrant_creative.pptx"),
        "bg_color": (249, 97, 103),         # coral
        "title_color": (255, 255, 255),
        "body_color": (255, 245, 200),      # warm cream
        "title_font": "Impact",
        "body_font": "Arial",
        "title_size": 42,
        "body_size": 17,
    },
    "nature_green": {
        "template": os.path.join(TEMPLATES_DIR, "nature_green.pptx"),
        "bg_color": (44, 95, 45),           # forest green
        "title_color": (245, 245, 245),
        "body_color": (200, 230, 200),
        "title_font": "Trebuchet MS",
        "body_font": "Calibri",
        "title_size": 38,
        "body_size": 16,
    },
    "berlin": {
    "template": os.path.join(TEMPLATES_DIR, "berlin.pptx"),
    "bg_color": (20, 20, 30),
    "title_color": (255, 200, 0),
    "body_color": (230, 230, 230),
    "title_font": "Arial Black",
    "body_font": "Calibri",
    "title_size": 40,
    "body_size": 17,
},
"slate_dark": {
    "template": os.path.join(TEMPLATES_DIR, "slate_dark.pptx"),
    "bg_color": (54, 69, 79),
    "title_color": (242, 242, 242),
    "body_color": (200, 210, 215),
    "title_font": "Trebuchet MS",
    "body_font": "Calibri Light",
    "title_size": 38,
    "body_size": 16,
},
"warm_terracotta": {
    "template": os.path.join(TEMPLATES_DIR, "warm_terracotta.pptx"),
    "bg_color": (184, 80, 66),
    "title_color": (255, 245, 220),
    "body_color": (231, 232, 209),
    "title_font": "Georgia",
    "body_font": "Calibri",
    "title_size": 38,
    "body_size": 16,
},
"ocean_depth": {
    "template": os.path.join(TEMPLATES_DIR, "ocean_depth.pptx"),
    "bg_color": (6, 90, 130),
    "title_color": (255, 255, 255),
    "body_color": (180, 230, 245),
    "title_font": "Calibri",
    "body_font": "Calibri Light",
    "title_size": 38,
    "body_size": 17,
},
"berry_bold": {
    "template": os.path.join(TEMPLATES_DIR, "berry_bold.pptx"),
    "bg_color": (109, 46, 70),
    "title_color": (236, 226, 208),
    "body_color": (210, 180, 190),
    "title_font": "Georgia",
    "body_font": "Calibri",
    "title_size": 38,
    "body_size": 16,
},
"coral_energy": {
    "template": os.path.join(TEMPLATES_DIR, "coral_energy.pptx"),
    "bg_color": (249, 97, 103),
    "title_color": (255, 245, 150),
    "body_color": (255, 255, 255),
    "title_font": "Impact",
    "body_font": "Arial",
    "title_size": 42,
    "body_size": 17,
},
"midnight_exec": {
    "template": os.path.join(TEMPLATES_DIR, "midnight_exec.pptx"),
    "bg_color": (30, 39, 97),
    "title_color": (202, 220, 252),
    "body_color": (255, 255, 255),
    "title_font": "Cambria",
    "body_font": "Calibri",
    "title_size": 40,
    "body_size": 17,
},
"vintage_wood": {
    "template": os.path.join(TEMPLATES_DIR, "vintage_wood.pptx"),
    "bg_color": (101, 67, 33),
    "title_color": (255, 240, 200),
    "body_color": (230, 210, 175),
    "title_font": "Palatino Linotype",
    "body_font": "Garamond",
    "title_size": 38,
    "body_size": 16,
},
"circuit_board": {
    "template": os.path.join(TEMPLATES_DIR, "circuit_board.pptx"),
    "bg_color": (10, 20, 10),
    "title_color": (0, 255, 100),
    "body_color": (180, 255, 180),
    "title_font": "Consolas",
    "body_font": "Courier New",
    "title_size": 38,
    "body_size": 16,
},
"celestial": {
    "template": os.path.join(TEMPLATES_DIR, "celestial.pptx"),
    "bg_color": (10, 10, 50),
    "title_color": (180, 200, 255),
    "body_color": (210, 220, 255),
    "title_font": "Arial Black",
    "body_font": "Calibri Light",
    "title_size": 40,
    "body_size": 17,
},
"sage_calm": {
    "template": os.path.join(TEMPLATES_DIR, "sage_calm.pptx"),
    "bg_color": (132, 181, 159),
    "title_color": (30, 30, 30),
    "body_color": (40, 60, 50),
    "title_font": "Georgia",
    "body_font": "Calibri Light",
    "title_size": 38,
    "body_size": 16,
},
"cherry_bold": {
    "template": os.path.join(TEMPLATES_DIR, "cherry_bold.pptx"),
    "bg_color": (153, 0, 17),
    "title_color": (255, 255, 255),
    "body_color": (255, 220, 220),
    "title_font": "Arial Black",
    "body_font": "Arial",
    "title_size": 42,
    "body_size": 17,
},
"droplet_fresh": {
    "template": os.path.join(TEMPLATES_DIR, "droplet_fresh.pptx"),
    "bg_color": (2, 160, 150),
    "title_color": (255, 255, 255),
    "body_color": (220, 248, 245),
    "title_font": "Trebuchet MS",
    "body_font": "Calibri",
    "title_size": 38,
    "body_size": 16,
},
"golden_exec": {
    "template": os.path.join(TEMPLATES_DIR, "golden_exec.pptx"),
    "bg_color": (50, 30, 10),
    "title_color": (212, 175, 55),
    "body_color": (240, 220, 170),
    "title_font": "Cambria",
    "body_font": "Garamond",
    "title_size": 40,
    "body_size": 17,
},
"pastel_soft": {
    "template": os.path.join(TEMPLATES_DIR, "pastel_soft.pptx"),
    "bg_color": (255, 240, 245),
    "title_color": (180, 100, 130),
    "body_color": (100, 80, 100),
    "title_font": "Trebuchet MS",
    "body_font": "Calibri Light",
    "title_size": 36,
    "body_size": 15,
},
}

# Fallback if no template file exists — pure programmatic styling
def _apply_theme_styles(prs: Presentation, theme: dict):
    """Apply color + font overrides to ALL slides in the presentation."""
    for slide in prs.slides:
        # Background
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*theme["bg_color"])

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            is_title = shape == slide.shapes.title
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = theme["title_font"] if is_title else theme["body_font"]
                    run.font.color.rgb = RGBColor(
                        *theme["title_color"] if is_title else theme["body_color"]
                    )
                    run.font.size = Pt(theme["title_size"] if is_title else theme["body_size"])
                    run.font.bold = is_title
                para.alignment = PP_ALIGN.CENTER if is_title else PP_ALIGN.LEFT


def download_pexels_image(query: str, save_path: str):
    """Download first Pexels image for a query."""
    headers = {"Authorization": Config.PEXELS_API_KEY}
    url = f"https://api.pexels.com/v1/search?query={query}&per_page=1"
    try:
        resp = requests.get(url, headers=headers, timeout=10)
        if resp.status_code == 200:
            photos = resp.json().get("photos", [])
            if photos:
                img_url = photos[0]["src"]["large"]
                img_resp = requests.get(img_url, timeout=10)
                if img_resp.status_code == 200:
                    with open(save_path, "wb") as f:
                        f.write(img_resp.content)
                    return save_path
    except Exception:
        pass
    return None


# ─── TOOLS ─────────────────────────────────────────────────────────────────────

@mcp.tool()
def create_presentation(title: str, theme: str = "minimal_light") -> str:
    """
    Create a new presentation using the specified theme template.
    Falls back to a blank presentation if template file is missing.
    """
    try:
        theme_data = THEMES.get(theme, THEMES["minimal_light"])
        template_path = theme_data["template"]

        if os.path.exists(template_path):
            prs = Presentation(template_path)
            # Use first slide layout (title slide)
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
        else:
            # Fallback: blank presentation
            prs = Presentation()
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)

        # Set title slide content
        if slide.shapes.title:
            slide.shapes.title.text = title
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = "Generated by AgentSlides"

        # Apply theme styles
        _apply_theme_styles(prs, theme_data)

        # Store chosen theme name in a custom property so other tools can read it
        prs.core_properties.keywords = theme  # reuse keywords field to persist theme choice

        prs.save(TEMP_FILE)
        return f"Presentation created with theme '{theme}'"
    except Exception as e:
        return f"[ERROR in create_presentation] {str(e)}"


@mcp.tool()
def add_slide_with_image(title: str, bullets: List[str]) -> str:
    """Add a slide with title, bullet points, and a relevant Pexels image."""
    try:
        if not os.path.exists(TEMP_FILE):
            return "Error: Presentation not initialized."

        prs = Presentation(TEMP_FILE)

        # Recover theme from persisted keyword
        theme_name = prs.core_properties.keywords or "minimal_light"
        theme_data = THEMES.get(theme_name, THEMES["minimal_light"])

        layout_idx = 3 if len(prs.slide_layouts) > 3 else 1
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

        # Title
        if slide.shapes.title:
            slide.shapes.title.text = title

        # Bullets — find content placeholder
        left_ph, right_ph = None, None
        for shape in slide.placeholders:
            idx = shape.placeholder_format.idx
            if idx == 1:
                left_ph = shape
            elif idx == 2:
                right_ph = shape

        if left_ph:
            tf = left_ph.text_frame
            tf.clear()
            for i, point in enumerate(bullets):
                if i == 0:
                    tf.text = point
                    for run in tf.paragraphs[0].runs:
                        run.font.size = Pt(theme_data["body_size"])
                        run.font.color.rgb = RGBColor(*theme_data["body_color"])
                else:
                    p = tf.add_paragraph()
                    p.text = point
                    p.level = 0
                    for run in p.runs:
                        run.font.size = Pt(theme_data["body_size"])
                        run.font.color.rgb = RGBColor(*theme_data["body_color"])

        # Pexels image
        query = f"{title} {' '.join(b.split(':')[0] for b in bullets[:2])}"
        img_path = os.path.join(BASE_DIR, f"pexels_{title.replace(' ', '_')[:30]}.jpg")
        img_file = download_pexels_image(query, img_path)

        if img_file and os.path.exists(img_file):
            if right_ph:
                try:
                    slide.shapes.add_picture(
                        img_file,
                        right_ph.left, right_ph.top,
                        right_ph.width, right_ph.height
                    )
                except Exception:
                    pass
            else:
                try:
                    slide.shapes.add_picture(img_file, Emu(4500000), Emu(1500000), Emu(4000000), Emu(3000000))
                except Exception:
                    pass

        # Apply theme styling to new slide
        _apply_theme_styles(prs, theme_data)

        prs.save(TEMP_FILE)
        return f"Slide added: {title}"
    except Exception as e:
        return f"[ERROR in add_slide_with_image] {str(e)}"


@mcp.tool()
def add_slide(title: str, bullets: List[str]) -> str:
    """Add a plain slide (no image) with theme styling."""
    try:
        if not os.path.exists(TEMP_FILE):
            return "Error: Presentation not initialized."

        prs = Presentation(TEMP_FILE)
        theme_name = prs.core_properties.keywords or "minimal_light"
        theme_data = THEMES.get(theme_name, THEMES["minimal_light"])

        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title

        tf = slide.placeholders[1].text_frame
        tf.clear()
        for i, point in enumerate(bullets):
            if i == 0:
                tf.text = point
            else:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0

        _apply_theme_styles(prs, theme_data)
        prs.save(TEMP_FILE)
        return f"Slide added: {title}"
    except Exception as e:
        return f"[ERROR in add_slide] {str(e)}"


@mcp.tool()
def save_presentation(filename: str) -> str:
    """Save the final presentation to the output folder."""
    try:
        if not os.path.exists(TEMP_FILE):
            return "Error: No presentation to save."
        final_path = os.path.join(BASE_DIR, os.path.basename(filename))
        shutil.copy(TEMP_FILE, final_path)
        os.remove(TEMP_FILE)  # cleanup temp
        return f"Presentation saved at: {final_path}"
    except Exception as e:
        return f"[ERROR in save_presentation] {str(e)}"


if __name__ == "__main__":
    mcp.run()